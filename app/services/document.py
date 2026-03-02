"""Document processor: extract text from PPTX, PDF, DOCX, and TXT files.

For PowerPoint files, attempts to convert slides to PNG images via LibreOffice
(if installed) in addition to text extraction.
"""

from __future__ import annotations

import logging
import os
import subprocess
from pathlib import Path
from typing import List, Optional, Tuple

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pptx", ".ppt", ".pdf", ".docx", ".doc", ".txt"}


class DocumentProcessor:
    """Processes documents and extracts text content."""

    def process(self, file_path: str, filename: str) -> Tuple[str, List[str]]:
        """Extract text from *file_path*.

        Returns:
            (extracted_text, list_of_image_paths)
        """
        ext = Path(filename).suffix.lower()

        if ext in (".pptx", ".ppt"):
            return self._process_pptx(file_path)
        if ext == ".pdf":
            return self._process_pdf(file_path)
        if ext in (".docx", ".doc"):
            return self._process_docx(file_path)
        if ext == ".txt":
            return self._process_txt(file_path)

        raise ValueError(
            f"未対応のファイル形式です: {ext}。"
            f"対応形式: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    # ------------------------------------------------------------------
    # Private helpers
    # ------------------------------------------------------------------

    def _process_pptx(self, file_path: str) -> Tuple[str, List[str]]:
        """Extract text from PowerPoint and optionally convert slides to PNG."""
        try:
            from pptx import Presentation  # type: ignore
        except ImportError as exc:
            raise RuntimeError("python-pptx がインストールされていません") from exc

        prs = Presentation(file_path)
        slides_text: List[str] = []

        for i, slide in enumerate(prs.slides, 1):
            parts = [f"--- スライド {i} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
            slides_text.append("\n".join(parts))

        full_text = "\n\n".join(slides_text)

        # Optional: convert to PNG images using LibreOffice
        image_paths: List[str] = []
        try:
            image_paths = self._convert_to_images(file_path)
        except Exception as exc:  # pylint: disable=broad-except
            logger.warning("LibreOffice によるスライド変換に失敗しました: %s", exc)

        return full_text, image_paths

    def _convert_to_images(self, file_path: str) -> List[str]:
        """Convert a PPTX file to PNG images using LibreOffice."""
        libreoffice = self._find_libreoffice()
        if not libreoffice:
            return []

        output_dir = os.path.join(os.path.dirname(file_path), "slides")
        os.makedirs(output_dir, exist_ok=True)

        result = subprocess.run(
            [libreoffice, "--headless", "--convert-to", "png",
             "--outdir", output_dir, file_path],
            capture_output=True,
            text=True,
            timeout=120,
        )

        if result.returncode != 0:
            logger.warning("LibreOffice 変換失敗: %s", result.stderr)
            return []

        images = sorted(Path(output_dir).glob("*.png"))
        return [str(p) for p in images]

    @staticmethod
    def _find_libreoffice() -> Optional[str]:
        """Return the path to LibreOffice / soffice, or None if not found."""
        candidates = [
            "libreoffice",
            "soffice",
            "/usr/bin/libreoffice",
            "/usr/bin/soffice",
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        ]
        for cmd in candidates:
            try:
                r = subprocess.run(
                    [cmd, "--version"],
                    capture_output=True,
                    timeout=5,
                )
                if r.returncode == 0:
                    return cmd
            except (FileNotFoundError, subprocess.TimeoutExpired):
                continue
        return None

    def _process_pdf(self, file_path: str) -> Tuple[str, List[str]]:
        """Extract text from a PDF file."""
        try:
            import pdfplumber  # type: ignore

            text_parts: List[str] = []
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"--- ページ {i} ---\n{page_text}")
            return "\n\n".join(text_parts), []

        except ImportError:
            logger.warning("pdfplumber が利用できません。PyPDF2 にフォールバックします")

        # Fallback: PyPDF2
        try:
            import PyPDF2  # type: ignore

            text_parts = []
            with open(file_path, "rb") as fh:
                reader = PyPDF2.PdfReader(fh)
                for i, page in enumerate(reader.pages, 1):
                    text = page.extract_text()
                    if text:
                        text_parts.append(f"--- ページ {i} ---\n{text}")
            return "\n\n".join(text_parts), []

        except ImportError as exc:
            raise RuntimeError(
                "pdfplumber または PyPDF2 がインストールされていません"
            ) from exc

    def _process_docx(self, file_path: str) -> Tuple[str, List[str]]:
        """Extract text from a Word document."""
        try:
            import docx  # type: ignore
        except ImportError as exc:
            raise RuntimeError("python-docx がインストールされていません") from exc

        doc = docx.Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs), []

    @staticmethod
    def _process_txt(file_path: str) -> Tuple[str, List[str]]:
        """Read a plain-text file."""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as fh:
            return fh.read(), []
